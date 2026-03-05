# backend/services/slides.py

import os
from pptx import Presentation
from pptx.util import Inches, Pt


class SlideService:
    def generate_pptx(self, slides_data: dict, output_path: str) -> str:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content

        for slide_info in slides_data.get("slides", []):
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = slide_info.get("title", "")
            tf = body.text_frame
            tf.clear()

            for point in slide_info.get("points", []):
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

        prs.save(output_path)
        return output_path
